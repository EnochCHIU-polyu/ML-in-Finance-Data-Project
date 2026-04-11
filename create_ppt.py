"""
Generate COMP5564 presentation (10 core + 1 backup slide) from project assets.
Run: .venv/bin/python create_ppt.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Paths
BASE = os.path.dirname(os.path.abspath(__file__))
PF = os.path.join(BASE, "project_folder")
REG_GRAPH = os.path.join(PF, "02_stock_price_regression", "graph")
CLU_GRAPH = os.path.join(PF, "03_stock_clustering_analysis", "graph")
OUT = os.path.join(PF, "COMP5564_PRESENTATION.pptx")

# Brand colors
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
CHARCOAL = RGBColor(0x2C, 0x2C, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x27, 0xAE, 0x60)
GOLD = RGBColor(0xD4, 0xA0, 0x1E)
SOFT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
SOFT_GOLD = RGBColor(0xFE, 0xF5, 0xE7)
SOFT_GREEN = RGBColor(0xE6, 0xF4, 0xEA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    color=CHARCOAL,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    return box


def add_bullets(slide, left, top, width, height, items, size=15, color=CHARCOAL):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "- " + item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(3)
    return box


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.isfile(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)
    else:
        add_textbox(
            slide,
            left,
            top,
            Inches(4),
            Inches(0.4),
            "[Missing image: {}]".format(os.path.basename(path)),
            size=11,
            color=ACCENT,
        )


def content_slide(slide, title):
    add_bg(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.95), NAVY)
    add_rect(slide, Inches(0), Inches(0.95), W, Inches(0.04), GOLD)
    add_textbox(slide, Inches(0.6), Inches(0.12), Inches(12.0), Inches(0.6), title, size=25, color=WHITE, bold=True)
    return Inches(1.2)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# 1) Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, Inches(0), Inches(5.8), W, Inches(0.06), GOLD)
add_rect(slide, Inches(0), Inches(6.0), W, Inches(1.5), CHARCOAL)
add_textbox(slide, Inches(0.8), Inches(1.4), Inches(12), Inches(1.4), "Machine Learning in Finance", 44, WHITE, True)
add_textbox(
    slide,
    Inches(0.8),
    Inches(2.8),
    Inches(12),
    Inches(0.8),
    "Cluster-Aware Sentiment and Multi-Meta Prediction for S&P 500",
    26,
    RGBColor(0xBB, 0xDD, 0xFF),
)
add_textbox(slide, Inches(0.8), Inches(4.6), Inches(8), Inches(0.5), "COMP5564 | CHIU Yee Lok (25012923G)", 18, WHITE)
add_textbox(slide, Inches(0.8), Inches(6.2), Inches(8), Inches(0.4), "April 2026", 15, RGBColor(0xAA, 0xAA, 0xAA))
set_notes(slide, "Open with the core thesis: two implementation designs, one system. Explain this is a deployment-focused ML in finance project.")

# 2) Problem, Goal, and Blueprint
slide = prs.slides.add_slide(prs.slide_layouts[6])
y = content_slide(slide, "Problem, Goal, and Presentation Blueprint")
add_bullets(
    slide,
    Inches(0.6),
    y,
    Inches(6.3),
    Inches(3.2),
    [
        "S&P 500 prediction is nonlinear, noisy, and regime-sensitive.",
        "Classical econometrics underfits mixed technical + news signals.",
        "Goal: improve predictive utility and tradable signal quality.",
        "Scope stock universe: 505 names; focus examples: AAPL, AMZN, GOOG, MSFT, NVDA.",
        "Strict temporal protocol and lagged features to avoid leakage.",
    ],
    size=16,
)
add_rect(slide, Inches(7.1), y, Inches(5.6), Inches(4.5), LIGHT)
add_textbox(slide, Inches(7.3), y + Inches(0.2), Inches(5.2), Inches(0.5), "Talk Flow", 18, NAVY, True)
add_bullets(
    slide,
    Inches(7.3),
    y + Inches(0.8),
    Inches(5.2),
    Inches(3.4),
    [
        "Dual-task architecture",
        "Implementation A: Cluster News design",
        "Implementation B: Multi-Meta Prediction design",
        "Results and robustness checks",
        "Limitations, controls, and deployment stance",
    ],
    size=14,
)
set_notes(slide, "State the business problem and what this deck will prove: not only prediction accuracy but portfolio-level utility.")

# 3) Dual-task architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
y = content_slide(slide, "Dual-Task Architecture: Regression x Clustering")
add_rect(slide, Inches(0.5), y, Inches(5.9), Inches(2.5), SOFT_BLUE)
add_textbox(slide, Inches(0.7), y + Inches(0.1), Inches(5.4), Inches(0.4), "Task 1: Stock Price Regression", 19, NAVY, True)
add_bullets(
    slide,
    Inches(0.7),
    y + Inches(0.6),
    Inches(5.4),
    Inches(1.8),
    [
        "Predict next-day close with XGBoost.",
        "Features: technicals, sentiment, volatility, momentum.",
        "Metrics: MAE, RMSE, R2, directional hit ratio.",
    ],
    size=14,
)
add_rect(slide, Inches(6.9), y, Inches(5.9), Inches(2.5), SOFT_GOLD)
add_textbox(slide, Inches(7.1), y + Inches(0.1), Inches(5.4), Inches(0.4), "Task 2: Stock Clustering", 19, GOLD, True)
add_bullets(
    slide,
    Inches(7.1),
    y + Inches(0.6),
    Inches(5.4),
    Inches(1.8),
    [
        "Cluster 505 stocks by behavior.",
        "K-Means (K=6) selected by composite score.",
        "Cluster IDs used as signal routing structure.",
    ],
    size=14,
)
add_rect(slide, Inches(2.8), y + Inches(2.85), Inches(7.8), Inches(0.06), ACCENT)
add_textbox(
    slide,
    Inches(1.8),
    y + Inches(3.1),
    Inches(9.8),
    Inches(1.0),
    "Coupling insight: clustering is not an appendix. It routes sentiment information into prediction.",
    17,
    NAVY,
    False,
    PP_ALIGN.CENTER,
)
set_notes(slide, "Explain the system-level idea: unsupervised structure feeds supervised learning through cluster-aware information routing.")

# 4) Implementation Design A
slide = prs.slides.add_slide(prs.slide_layouts[6])
y = content_slide(slide, "Implementation Design A: Cluster News Prediction")
add_bullets(
    slide,
    Inches(0.6),
    y,
    Inches(6.2),
    Inches(3.3),
    [
        "Problem: direct ticker news is sparse and uneven across names.",
        "Design: aggregate peer-cluster news when direct news is missing.",
        "Scoring: semantic similarity x recency decay x source reliability x sentiment polarity.",
        "Stacked residual design isolates sentiment contribution over price baseline.",
        "Grid-searched decay remains stable around L=10 and lambda=0.2.",
    ],
    size=15,
)
add_rect(slide, Inches(0.5), y + Inches(3.45), Inches(6.3), Inches(2.3), SOFT_GREEN)
add_textbox(slide, Inches(0.7), y + Inches(3.6), Inches(5.9), Inches(0.4), "Observed impact", 16, GREEN, True)
add_bullets(
    slide,
    Inches(0.7),
    y + Inches(4.1),
    Inches(5.9),
    Inches(1.5),
    [
        "AMZN MAE down about 47 percent; GOOG MAE down about 46 percent.",
        "Best stacked run: MAE 0.918, R2 0.999688, DHR 0.515.",
    ],
    size=13,
)
add_image_safe(slide, os.path.join(REG_GRAPH, "02D_prediction_vs_actual_close_GOOG.png"), Inches(7.0), y, width=Inches(5.8))
set_notes(slide, "Present Design A as the first core innovation. Highlight sparsity conversion: cluster peers provide signal when direct ticker news is absent.")

# 5) Implementation Design B
slide = prs.slides.add_slide(prs.slide_layouts[6])
y = content_slide(slide, "Implementation Design B: Multi-Meta Prediction")
add_bullets(
    slide,
    Inches(0.6),
    y,
    Inches(6.2),
    Inches(3.1),
    [
        "Problem: single-horizon T+1 forecasts are not portfolio-horizon aware.",
        "Design: combine T+1, T+5, T+10, T+15 heads through a learned meta layer.",
        "Meta controls: regime filter, direction agreement, velocity alignment.",
        "Ridge meta model prioritized for stability; XGB meta as secondary benchmark.",
        "Outcome: stronger risk-adjusted return with lower drawdown.",
    ],
    size=15,
)
add_rect(slide, Inches(7.0), y, Inches(5.8), Inches(4.4), LIGHT)
add_textbox(slide, Inches(7.2), y + Inches(0.2), Inches(5.4), Inches(0.4), "Performance snapshot", 16, NAVY, True)
add_bullets(
    slide,
    Inches(7.2),
    y + Inches(0.8),
    Inches(5.4),
    Inches(3.2),
    [
        "Learned Meta Ridge: Sharpe 8.10, Return +10.5 percent, MaxDD -0.8 percent, DHR 0.681.",
        "Learned Meta XGB: Sharpe 4.72, Return +6.5 percent.",
        "Buy-and-Hold baseline: Sharpe 6.09, Return +8.2 percent.",
    ],
    size=13,
)
add_textbox(
    slide,
    Inches(0.6),
    y + Inches(3.5),
    Inches(6.2),
    Inches(1.8),
    "Interpretation: long-horizon heads provide regime confidence to modulate short-horizon exposure.",
    14,
    NAVY,
    True,
)
set_notes(slide, "Present Design B as the second core innovation. Emphasize that horizon fusion changes position conviction, not just forecast error.")

# 6) Core empirical results
slide = prs.slides.add_slide(prs.slide_layouts[6])
y = content_slide(slide, "Core Empirical Results")
add_rect(slide, Inches(0.5), y, Inches(3.9), Inches(3.6), SOFT_BLUE)
add_textbox(slide, Inches(0.7), y + Inches(0.1), Inches(3.5), Inches(0.4), "Regression", 16, NAVY, True)
add_bullets(
    slide,
    Inches(0.7),
    y + Inches(0.6),
    Inches(3.5),
    Inches(2.7),
    [
        "XGBoost selected over deep models.",
        "Stacked MAE 0.918.",
        "R2 0.999688.",
        "DHR 0.515.",
    ],
    size=13,
)
add_rect(slide, Inches(4.8), y, Inches(3.9), Inches(3.6), SOFT_GOLD)
add_textbox(slide, Inches(5.0), y + Inches(0.1), Inches(3.5), Inches(0.4), "Clustering", 16, GOLD, True)
add_bullets(
    slide,
    Inches(5.0),
    y + Inches(0.6),
    Inches(3.5),
    Inches(2.7),
    [
        "K-Means chosen at K=6.",
        "Composite score 0.577.",
        "Silhouette 0.250.",
        "Zero noise ratio.",
    ],
    size=13,
)
add_rect(slide, Inches(9.1), y, Inches(3.9), Inches(3.6), SOFT_GREEN)
add_textbox(slide, Inches(9.3), y + Inches(0.1), Inches(3.5), Inches(0.4), "Trading", 16, GREEN, True)
add_bullets(
    slide,
    Inches(9.3),
    y + Inches(0.6),
    Inches(3.5),
    Inches(2.7),
    [
        "Meta Ridge Sharpe 8.10.",
        "Return +10.5 percent.",
        "Max drawdown -0.8 percent.",
        "Outperforms buy-and-hold.",
    ],
    size=13,
)
add_image_safe(slide, os.path.join(CLU_GRAPH, "03A_02_pca_scatter.png"), Inches(0.5), y + Inches(3.9), width=Inches(6.1))
add_image_safe(slide, os.path.join(REG_GRAPH, "03_regression_metrics_comparison.png"), Inches(6.8), y + Inches(3.9), width=Inches(6.0))
set_notes(slide, "Summarize performance on three layers: prediction, structure, and trading. Keep this slide quantitative and fast.")

# 7) Trading behavior and validation
slide = prs.slides.add_slide(prs.slide_layouts[6])
y = content_slide(slide, "Trading Behavior, WFO Validation, and Alpha Diagnostics")
add_image_safe(slide, os.path.join(REG_GRAPH, "02F_01_cumulative_returns.png"), Inches(0.5), y, width=Inches(6.0))
add_image_safe(slide, os.path.join(REG_GRAPH, "02F_02_drawdown.png"), Inches(6.8), y, width=Inches(6.0))
add_rect(slide, Inches(0.5), y + Inches(4.0), Inches(12.3), Inches(1.8), LIGHT)
add_bullets(
    slide,
    Inches(0.7),
    y + Inches(4.15),
    Inches(11.9),
    Inches(1.5),
    [
        "WFO comparison added to prevent fixed-split optimism in final reporting.",
        "IC and ICIR diagnostics included for cross-sectional stock-selection signal quality.",
        "Primary interpretation remains conservative until larger and multi-regime OOS confirms stability.",
    ],
    size=13,
)
set_notes(slide, "Use this as the institutional validation slide: equity behavior, drawdown control, WFO realism, and IC/ICIR signal diagnostics.")

# 8) Controls and risk governance
slide = prs.slides.add_slide(prs.slide_layouts[6])
y = content_slide(slide, "Controls, Auditability, and Governance")
add_rect(slide, Inches(0.5), y, Inches(6.1), Inches(5.6), SOFT_BLUE)
add_textbox(slide, Inches(0.7), y + Inches(0.2), Inches(5.7), Inches(0.4), "Implemented controls", 17, NAVY, True)
add_bullets(
    slide,
    Inches(0.7),
    y + Inches(0.8),
    Inches(5.7),
    Inches(4.7),
    [
        "Temporal split and lag enforcement to prevent look-ahead leakage.",
        "Model comparison under shared protocol, not cherry-picked runs.",
        "Feature and strategy diagnostics documented in notebook outputs.",
        "WFO and uncertainty framing integrated into report narrative.",
        "Benchmark comparison includes passive baseline and risk metrics.",
    ],
    size=14,
)
add_rect(slide, Inches(6.9), y, Inches(5.9), Inches(5.6), SOFT_GOLD)
add_textbox(slide, Inches(7.1), y + Inches(0.2), Inches(5.5), Inches(0.4), "Remaining risks", 17, GOLD, True)
add_bullets(
    slide,
    Inches(7.1),
    y + Inches(0.8),
    Inches(5.5),
    Inches(4.7),
    [
        "No explicit slippage and impact model in current backtest.",
        "Static cluster topology may drift across regimes.",
        "Current OOS horizon is still limited for production confidence.",
        "Headline coverage quality varies by source and stock.",
        "Further stress testing needed for bear and shock windows.",
    ],
    size=14,
)
set_notes(slide, "Balance strengths and open risks. This slide demonstrates institutional-grade skepticism and control awareness.")

# 9) Limitations and next iteration
slide = prs.slides.add_slide(prs.slide_layouts[6])
y = content_slide(slide, "Limitations and Next Iteration Plan")
add_bullets(
    slide,
    Inches(0.6),
    y,
    Inches(12.2),
    Inches(2.7),
    [
        "Current portfolio simulation excludes full transaction cost stack and market impact.",
        "Cluster structure is static; no online updating mechanism yet.",
        "Model confidence is point-estimate heavy; interval prediction is limited.",
        "Validation breadth should expand with longer OOS and additional market regimes.",
    ],
    size=16,
)
add_rect(slide, Inches(0.5), y + Inches(2.9), Inches(12.3), Inches(2.8), SOFT_GREEN)
add_textbox(slide, Inches(0.7), y + Inches(3.1), Inches(11.9), Inches(0.4), "Next iteration priorities", 17, GREEN, True)
add_bullets(
    slide,
    Inches(0.7),
    y + Inches(3.7),
    Inches(11.9),
    Inches(1.8),
    [
        "Dynamic clustering with rolling updates and drift monitoring.",
        "Cost-aware backtesting with slippage, commissions, and liquidity penalties.",
        "Uncertainty-aware position sizing using prediction intervals.",
        "Regime-conditioned retraining cadence and model governance checkpoints.",
    ],
    size=14,
)
set_notes(slide, "Present a credible roadmap. Prioritize what would be done before any real capital deployment.")

# 10) Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11.8), Inches(0.8), "Conclusion", 36, WHITE, True)
add_rect(slide, Inches(0.8), Inches(1.75), Inches(3.0), Inches(0.04), GOLD)
add_bullets(
    slide,
    Inches(0.8),
    Inches(2.1),
    Inches(11.8),
    Inches(3.6),
    [
        "Design A (Cluster News) improves information coverage under sparse news conditions.",
        "Design B (Multi-Meta Prediction) improves risk-adjusted trading outcomes.",
        "Combined system delivers gains in both predictive and deployment-facing metrics.",
        "Evidence is promising, but deployment recommendation remains controlled and staged.",
    ],
    size=19,
    color=WHITE,
)
add_rect(slide, Inches(0.8), Inches(5.7), Inches(11.8), Inches(1.1), GOLD)
add_textbox(
    slide,
    Inches(1.0),
    Inches(5.95),
    Inches(11.4),
    Inches(0.8),
    "Core thesis: clustering serves as a structural signal router, not only an unsupervised analysis block.",
    16,
    NAVY,
    True,
    PP_ALIGN.CENTER,
)
set_notes(slide, "Close with one message: architecture-level coupling of clustering and forecasting is the source of edge, but deployment must remain risk-controlled.")

# 11) Backup / Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CHARCOAL)
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(1.1), "Thank You - Q&A", 44, WHITE, True, PP_ALIGN.CENTER)
add_rect(slide, Inches(5.0), Inches(3.3), Inches(3.2), Inches(0.04), GOLD)
add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.8), Inches(0.5), "Backup metrics for discussion", 18, GOLD, True, PP_ALIGN.CENTER)
add_bullets(
    slide,
    Inches(1.2),
    Inches(4.7),
    Inches(10.9),
    Inches(1.8),
    [
        "Regression: MAE 0.918, R2 0.999688, DHR 0.515.",
        "Clustering: K-Means K=6, composite 0.577, zero noise ratio.",
        "Trading: Meta Ridge Sharpe 8.10, return +10.5 percent, max drawdown -0.8 percent.",
    ],
    size=14,
    color=RGBColor(0xCC, 0xCC, 0xCC),
)
set_notes(slide, "Use this slide to answer detail questions on assumptions, metrics, and robustness checks.")


# Save
prs.save(OUT)
print("Saved:", OUT)
print("Slides generated:", len(prs.slides))
